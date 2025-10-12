import os
import io
import re
import streamlit as st
import requests
from io import StringIO
import docx
import PyPDF2
import pandas as pd

# -----------------------------
# Optional Packages
# -----------------------------
try:
    import docx
    docx_available = True
except ImportError:
    docx_available = False
    st.warning("python-docx not installed. DOCX files will not be supported.")

try:
    from pptx import Presentation
    pptx_available = True
except ImportError:
    pptx_available = False
    st.warning("python-pptx not installed. PPTX files will not be supported.")

# -----------------------------
# Session State
# -----------------------------
if "notes" not in st.session_state:
    st.session_state.notes = {}
if "refresh_notes" not in st.session_state:
    st.session_state.refresh_notes = False


# -----------------------------
# Session State
# -----------------------------
if "notes" not in st.session_state:
    st.session_state.notes = {}
if "refresh_notes" not in st.session_state:
    st.session_state.refresh_notes = False

# -----------------------------
# API Key
# -----------------------------
API_KEY = st.secrets["my_secrets"]["api_key"]
if not API_KEY:
    st.error("API key not found! Please check your Streamlit secrets.")
    st.stop()

# -----------------------------
# Page Config
# -----------------------------
st.set_page_config(page_title="AI Study Buddy", layout="wide", page_icon="🧠")

# -----------------------------
# Sidebar
# -----------------------------
st.sidebar.markdown("## 🧠 AI Study Buddy")
st.sidebar.write("Welcome! Use this tool to:")
st.sidebar.markdown("""
- **Explain** your notes in simple terms  
- **Summarize** long documents quickly  
- **Generate Quizzes** for self-assessment  
- **Create Flashcards**  
- **Save Notes**
""")

difficulty = st.sidebar.selectbox("AI Difficulty Level", ["Beginner","Intermediate","Advanced"])

# Notes
st.sidebar.markdown("### 🗒️ Notes")
note_title = st.sidebar.text_input("Note Title")
note_content = st.sidebar.text_area("Write your note here...", height=150)
col1, col2 = st.sidebar.columns(2)
if col1.button("💾 Save Note") and note_title.strip() and note_content.strip():
    st.session_state.notes[note_title] = note_content
    st.success("Note saved!")
if col2.button("🔄 Refresh Notes"):
    st.session_state.notes = {}
    st.session_state.refresh_notes = True
if st.session_state.refresh_notes:
    st.session_state.refresh_notes = False
    st.sidebar.info("Notes cleared! Start a new note now.")
if st.session_state.notes:
    st.sidebar.markdown("#### Saved Notes")
    for title, content in st.session_state.notes.items():
        st.sidebar.write(f"**{title}**")
        st.sidebar.write(content[:200]+"..." if len(content)>200 else content)
    combined_notes=""
    for title, content in st.session_state.notes.items():
        combined_notes += f"{title}\n{'-'*len(title)}\n{content}\n\n"
    st.sidebar.download_button("Download as .txt", combined_notes, "AI_Study_Notes.txt")
    try:
        from docx import Document
        doc = Document()
        for title, content in st.session_state.notes.items():
            doc.add_heading(title, level=2)
            doc.add_paragraph(content)
        doc_stream = io.BytesIO()
        doc.save(doc_stream)
        doc_stream.seek(0)
        st.sidebar.download_button("Download as .docx", doc_stream, "AI_Study_Notes.docx")
    except:
        pass

# -----------------------------
# App Title
# -----------------------------
st.markdown('<p class="big-font">📚Study Buddy</p>', unsafe_allow_html=True)
st.write("Upload notes or type text, then select a task below!")

# -----------------------------
# User Input
# -----------------------------
option = st.radio("Choose input type:", ("Text Input","Upload File"))
user_content = ""

if option=="Text Input":
    user_content = st.text_area("Enter your text/question here:", height=180)
elif option=="Upload File":
    uploaded_file = st.file_uploader("Upload a file (txt,pdf,docx,pptx)", type=["txt","pdf","docx","pptx"])
    if uploaded_file:
        file_type = uploaded_file.name.split('.')[-1].lower()
        try:
            if file_type=="txt":
                user_content = StringIO(uploaded_file.getvalue().decode("utf-8")).read()
            elif file_type=="pdf":
                pdf_reader = PyPDF2.PdfReader(uploaded_file)
                user_content=""
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text: user_content += text + "\n"
            elif file_type=="docx":
                doc = docx.Document(uploaded_file)
                user_content="\n".join([p.text for p in doc.paragraphs])
            elif file_type=="pptx" and pptx_available:
                prs = Presentation(uploaded_file)
                user_content=""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape,"text"):
                            user_content += shape.text + "\n"
            st.success("File content loaded successfully!")
        except Exception as e:
            st.error(f"Failed to read file: {e}")

# -----------------------------
# Tabs
# -----------------------------
tabs = st.tabs(["📝 Explain","📑 Summarize","❓ Quiz","🃏 Flashcards"])
task_names = ["Explain","Summarize","Quiz","Flashcards"]

# -----------------------------
# AI Request
# -----------------------------
def get_ai_response(prompt):
    url = "https://api.perplexity.ai/chat/completions"
    headers = {"Authorization":f"Bearer {API_KEY}", "Content-Type":"application/json"}
    data = {"model":"sonar-pro","messages":[{"role":"user","content":prompt}]}
    try:
        response = requests.post(url, headers=headers, json=data)
        if response.status_code==200:
            result = response.json()
            content = result.get("choices",[{}])[0].get("message",{}).get("content","")
            cleaned_content = re.sub(r'\[\d+\]','',content)
            sources = result.get("citations",[]) or []
            return cleaned_content.strip(), sources
        else:
            return f"Error {response.status_code}: {response.text}", []
    except Exception as e:
        return f"Request failed: {e}", []

# -----------------------------
# Generate AI Output with citations
# -----------------------------
for idx, tab in enumerate(tabs):
    with tab:
        task_name = task_names[idx]
        st.write(f"### {task_name} Task")
        if st.button(f"Generate {task_name} Output", key=idx) and user_content.strip():
            with st.spinner("Generating AI output..."):
                prompt=""
                if task_name=="Explain":
                    prompt=f"Explain the following content simply ({difficulty} level):\n{user_content}"
                elif task_name=="Summarize":
                    prompt=f"Summarize the following content in bullet points ({difficulty} level):\n{user_content}"
                elif task_name=="Quiz":
                    prompt=f"Create 5-10 multiple choice or short answer questions from the following content:\n{user_content}"
                elif task_name=="Flashcards":
                    prompt=f"Create numbered flashcards (Q1/A1 format) from the following content:\n{user_content}"

                ai_output, sources = get_ai_response(prompt)

                # --- Explain ---
                if task_name=="Explain":
                    content_with_citations = ai_output
                    if sources:
                        for idx_src, src in enumerate(sources, start=1):
                            content_with_citations += f" [{idx_src}]"

                    st.markdown(
                        f'<div class="card"><div class="card-answer">{content_with_citations}</div></div>',
                        unsafe_allow_html=True
                    )

                    if sources:
                        st.markdown("### 🔗 Sources")
                        for idx_src, src in enumerate(sources, start=1):
                            st.markdown(f"{idx_src}. [{src}]({src})")

                # --- Summarize ---
                elif task_name=="Summarize":
                    points = [p.strip('-').strip() for p in ai_output.split("\n") if p.strip()]
                    points_html = ""
                    for i, p in enumerate(points, start=1):
                        if sources and i <= len(sources):
                            p += f" [{i}]"
                        points_html += f"<li style='margin-bottom:12px;'>{p}</li>"

                    st.markdown(
                        f'<div class="card"><ul style="padding-left:20px; line-height:1.8;">{points_html}</ul></div>',
                        unsafe_allow_html=True
                    )

                    if sources:
                        st.markdown("### 🔗 Sources")
                        for idx_src, src in enumerate(sources, start=1):
                            st.markdown(f"{idx_src}. [{src}]({src})")

                # --- Quiz with sources ---
                elif task_name=="Quiz":
                    quiz_items = re.split(r'\n(?=Q\d*:|Q\d+\.|Question \d+:)', ai_output)
                    for q_idx, item in enumerate(quiz_items, start=1):
                        if not item.strip(): 
                            continue
                        if "Answer:" in item:
                            parts = item.split("Answer:", 1)
                        elif "A:" in item:
                            parts = item.split("A:", 1)
                        else:
                            parts = item.strip().split("\n", 1)
                        q = parts[0].strip()
                        a = parts[1].strip() if len(parts) > 1 else "Answer not provided."

                        
                        st.markdown(
                            f"""
                            <div style='
                                border-left:4px solid #4a90e2; 
                                padding:15px; 
                                margin-bottom:10px;
                                border-radius:6px;
                            '>
                                <div style='font-weight:bold; font-size:16px;'>❓ Q{q_idx}: {q}</div>
                                <div style='margin-top:8px;'>💡 Answer: {a}</div>
                            </div>
                            """,
                            unsafe_allow_html=True
                        )

                    # --- Quiz with sources and styled answers ---
                elif task_name == "Quiz":
                # Split AI output into individual quiz items using regex
                    quiz_items = re.split(r'\n(?=Q\d*:|Q\d+\.|Question \d+:)', ai_output)

                    for q_idx, item in enumerate(quiz_items, start=1):
                        if not item.strip():
                            continue

                        # Extract question and answer
                        if "Answer:" in item:
                            parts = item.split("Answer:", 1)
                        elif "A:" in item:
                            parts = item.split("A:", 1)
                        else:
                            # If format is unexpected, split by first newline
                            parts = item.strip().split("\n", 1)

                        question = parts[0].strip()
                        answer = parts[1].strip() if len(parts) > 1 else "Answer not provided."

                        # Display question and answer in a card with styling
                        st.markdown(
                            f"""
                            <div style='
                                border-left:4px solid #4a90e2; 
                                padding:15px; 
                                margin-bottom:12px;
                                border-radius:6px;
                                background-color:#f9f9f9;
                            '>
                                <div style='font-weight:bold; font-size:16px;'>❓ Q{q_idx}: {question}</div>
                                <div style='margin-top:8px; font-size:15px;'>💡 Answer: {answer}</div>
                            </div>
                            """,
                            unsafe_allow_html=True
                            )


                # --- Flashcards ---
                elif task_name=="Flashcards":
                    flashcards = re.findall(r"(Q\d+:.*?\nA\d+:.*?)(?=Q\d+:|$)", ai_output, re.DOTALL)
                    for idx_fc, card in enumerate(flashcards, start=1):
                        q_match = re.search(r"Q\d+:\s*(.*)", card)
                        a_match = re.search(r"A\d+:\s*(.*)", card)
                        question = q_match.group(1).strip() if q_match else "Question not provided"
                        answer = a_match.group(1).strip() if a_match else "Answer not provided"
                        st.markdown(
                            f'<div class="card"><div class="card-question">❓ Flashcard {idx_fc}: {question}</div><div class="card-answer">💡 {answer}</div></div>',
                            unsafe_allow_html=True
                        )
